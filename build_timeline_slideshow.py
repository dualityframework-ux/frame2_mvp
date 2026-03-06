"""
build_timeline_slideshow.py
Generates a PPTX + standalone HTML slideshow from the Red Sox franchise timeline.

Slides produced
───────────────
  0  Cover — "Red Sox Franchise Timeline 1967–2024"
  1  The Story in Numbers — franchise KPI cards
  2  The Seven Eras — era map overview
  3  Era deep-dives × 7 (one per macro-era) — bar chart + stats table + standout season
 10  The Six Championships — WS year cards
 11  Process vs Result explained — scatter chart
 12  Wins Timeline — full 1967-2024 line/bar chart
 13  Era Radar comparison
 14  Key Players by Era
 15  Current State & What's Next
 16  Closing card
"""

import sys, os, io, math, textwrap
sys.path.insert(0, os.path.dirname(__file__))

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.lines import Line2D
from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as pns
from lxml import etree

from adapters.redsox_history_loader import load_redsox_history
from intelligence.redsox_history_engine import classify_process_vs_result

# ══════════════════════════════════════════════════════════════════════════════
# PALETTE & CONSTANTS
# ══════════════════════════════════════════════════════════════════════════════
RED    = RGBColor(0xC8, 0x10, 0x2E)
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
GOLD   = RGBColor(0xE4, 0xB8, 0x4D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MGRAY  = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
GREEN  = RGBColor(0x2D, 0xA8, 0x70)

# 16:9 slide dimensions
SW = Inches(13.33)
SH = Inches(7.5)

MACRO_ERAS = {
    "⚾ Yaz Era":              (list(range(1967, 1977)), "#4a90d9"),
    "💪 Rice & Lynn":           (list(range(1977, 1987)), "#e4b84d"),
    "🎯 Boggs & Clemens":       (list(range(1987, 1997)), "#7eb8f7"),
    "🔥 Pedro & Manny":         (list(range(1997, 2007)), "#c8102e"),
    "🏆 Francona Dynasty":      (list(range(2007, 2014)), "#2da870"),
    "🌟 Betts Emergence":       (list(range(2014, 2020)), "#c8940a"),
    "🔄 Rebuild & Now":         (list(range(2020, 2025)), "#999999"),
}

ERA_YEARS = {k: v[0] for k, v in MACRO_ERAS.items()}
ERA_COLORS_HEX = {k: v[1] for k, v in MACRO_ERAS.items()}

ERA_TAGLINES = {
    "⚾ Yaz Era":          "The Impossible Dream and the years that followed",
    "💪 Rice & Lynn":       "Power, heartbreak, and the '78 collapse",
    "🎯 Boggs & Clemens":   "Consistent talent; inconsistent October",
    "🔥 Pedro & Manny":     "The curse is reversed — finally",
    "🏆 Francona Dynasty":  "Back-to-back titles and the franchise's peak process era",
    "🌟 Betts Emergence":   "2018: the greatest regular season in team history",
    "🔄 Rebuild & Now":     "Post-Betts transition — the foundation is being set",
}

ERA_KEY_PLAYERS = {
    "⚾ Yaz Era":          "Carl Yastrzemski · Jim Lonborg · Carlton Fisk · Luis Tiant · Fred Lynn",
    "💪 Rice & Lynn":       "Jim Rice · Fred Lynn · Carl Yastrzemski · Dennis Eckersley · Carlton Fisk",
    "🎯 Boggs & Clemens":   "Roger Clemens · Wade Boggs · Dwight Evans · Mike Greenwell · Mo Vaughn",
    "🔥 Pedro & Manny":     "Pedro Martinez · Manny Ramirez · David Ortiz · Nomar Garciaparra · Curt Schilling",
    "🏆 Francona Dynasty":  "Josh Beckett · Jon Lester · David Ortiz · Dustin Pedroia · Jacoby Ellsbury",
    "🌟 Betts Emergence":   "Mookie Betts · Xander Bogaerts · Chris Sale · Rafael Devers · J.D. Martinez",
    "🔄 Rebuild & Now":     "Rafael Devers · Jarren Duran · Brayan Bello · Tanner Houck · Garrett Whitlock",
}

WS_SEASONS = [1967, 1975, 1986, 2004, 2007, 2013, 2018]

TAG_TEXT = {
    "signal":                "✅ Signal — process backed the result",
    "balanced":              "➡️ Balanced — result matched the process",
    "overperformed process": "🍀 Overperformed — won more than predicted",
    "underperformed process":"📉 Underperformed — won less than predicted",
}

# ══════════════════════════════════════════════════════════════════════════════
# CHART HELPERS
# ══════════════════════════════════════════════════════════════════════════════
def hex_to_rgb_f(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16)/255 for i in (0, 2, 4))

def fig_bytes(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def styled_ax(ax, title="", xlabel="", ylabel=""):
    ax.set_facecolor("#f8f8f8")
    ax.set_title(title, fontsize=11, fontweight="bold", color="#0d1b2a", pad=8)
    if xlabel: ax.set_xlabel(xlabel, fontsize=9, color="#555")
    if ylabel: ax.set_ylabel(ylabel, fontsize=9, color="#555")
    ax.tick_params(colors="#555", labelsize=8)
    for sp in ax.spines.values(): sp.set_color("#ccc")
    ax.grid(axis="y", color="#ddd", linestyle="--", linewidth=0.5)
    ax.set_axisbelow(True)

# Full-width wins bar for era slide
def era_wins_chart(sub_df, era_name, color_hex):
    fig, ax = plt.subplots(figsize=(9, 3.2), facecolor="white")
    seasons = sub_df["season"].astype(int).tolist()
    wins    = sub_df["wins"].tolist()
    c = hex_to_rgb_f(color_hex)
    bars = ax.bar(seasons, wins, color=c, width=0.7, zorder=3, edgecolor="white")
    ax.axhline(81, color="#c8102e", linestyle="--", linewidth=1.2, alpha=0.7, label=".500 line")
    for bar, w, s in zip(bars, wins, seasons):
        if sub_df[sub_df["season"]==s]["playoff_result"].values[0] == "world series":
            ax.text(bar.get_x()+bar.get_width()/2, w+1.2, "★",
                    ha="center", va="bottom", fontsize=13, color="#e4b84d")
        ax.text(bar.get_x()+bar.get_width()/2, w+1.5 if sub_df[sub_df["season"]==s]["playoff_result"].values[0]!="world series" else w+3,
                str(w), ha="center", va="bottom", fontsize=7.5, color="#0d1b2a")
    ax.set_xticks(seasons)
    ax.set_xticklabels([str(s) for s in seasons], fontsize=8, rotation=45)
    ax.set_ylim(50, 115)
    ax.legend(fontsize=8, loc="lower right")
    styled_ax(ax, f"Wins per Season — {era_name.split('(')[0].strip()}", ylabel="Wins")
    fig.tight_layout(pad=1.2)
    return fig_bytes(fig)

# Full franchise wins timeline
def franchise_timeline_chart(df):
    era_color_map = {}
    for era, (yrs, clr) in MACRO_ERAS.items():
        for y in yrs: era_color_map[y] = clr

    fig, ax = plt.subplots(figsize=(13, 4), facecolor="white")
    ax.set_facecolor("#f8f8f8")
    for _, row in df.iterrows():
        c = hex_to_rgb_f(era_color_map.get(int(row["season"]), "#aaa"))
        ax.bar(row["season"], row["wins"], color=c, width=0.75, zorder=3, edgecolor="white")
    ax.axhline(81, color="#c8102e", linestyle="--", linewidth=1.2, alpha=0.7, label=".500 (81 W)")

    ws = df[df["playoff_result"]=="world series"]
    ax.scatter(ws["season"], ws["wins"], s=140, color="#e4b84d", zorder=5, marker="*", label="World Series Win")

    # Era boundary lines
    for boundary in [1977, 1987, 1997, 2007, 2014, 2020]:
        ax.axvline(boundary-0.5, color="#aaa", linestyle=":", linewidth=0.8)

    ax.set_xlim(1965.5, 2025.5)
    ax.set_ylim(40, 125)
    ax.set_xlabel("Season", fontsize=9, color="#555")
    ax.set_ylabel("Wins", fontsize=9, color="#555")
    ax.set_title("Red Sox Win Totals — Every Season 1967–2024", fontsize=12,
                 fontweight="bold", color="#0d1b2a", pad=8)
    ax.tick_params(colors="#555", labelsize=8)
    for sp in ax.spines.values(): sp.set_color("#ccc")
    ax.grid(axis="y", color="#ddd", linestyle="--", linewidth=0.5)
    ax.set_axisbelow(True)

    legend_patches = [mpatches.Patch(color=hex_to_rgb_f(clr), label=era.split("(")[0].strip())
                      for era, (_, clr) in MACRO_ERAS.items()]
    legend_patches += [Line2D([0],[0], marker="*", color="w", markerfacecolor="#e4b84d",
                               markersize=10, label="World Series Win")]
    ax.legend(handles=legend_patches, fontsize=7.5, loc="upper left", ncol=2, frameon=True)
    fig.tight_layout(pad=1.2)
    return fig_bytes(fig)

# Scatter: wins vs run_diff
def scatter_chart(df):
    era_color_map = {}
    for era, (yrs, clr) in MACRO_ERAS.items():
        for y in yrs: era_color_map[y] = clr

    fig, ax = plt.subplots(figsize=(9, 5), facecolor="white")
    ax.set_facecolor("#f8f8f8")
    for _, row in df.iterrows():
        c = hex_to_rgb_f(era_color_map.get(int(row["season"]), "#aaa"))
        ax.scatter(row["run_diff"], row["wins"], color=c, s=55, alpha=0.85,
                   edgecolors="white", linewidth=0.4, zorder=3)
        if row["playoff_result"] == "world series":
            ax.scatter(row["run_diff"], row["wins"], s=130, facecolors="none",
                       edgecolors="#e4b84d", linewidth=2, zorder=4)
            ax.annotate(str(int(row["season"])), (row["run_diff"], row["wins"]),
                        textcoords="offset points", xytext=(5,3),
                        fontsize=7.5, color="#e4b84d", fontweight="bold")
    ax.axhline(81, color="#c8102e", linestyle="--", linewidth=1, alpha=0.6, label=".500 line")
    ax.axvline(0,  color="#c8102e", linestyle="--", linewidth=1, alpha=0.6)
    styled_ax(ax, "Process vs Result — Wins vs Run Differential",
              xlabel="Run Differential (process proxy)", ylabel="Wins")
    ax.text(-290, 82, ".500", fontsize=7.5, color="#c8102e", alpha=0.8)
    ax.text(2, 43, "Zero RD", fontsize=7.5, color="#c8102e", alpha=0.8, rotation=90)

    legend_patches = [mpatches.Patch(color=hex_to_rgb_f(clr), label=era.split("(")[0].strip())
                      for era, (_, clr) in MACRO_ERAS.items()]
    legend_patches += [Line2D([0],[0], marker="o", color="w", markerfacecolor="none",
                               markeredgecolor="#e4b84d", markeredgewidth=2, markersize=9,
                               label="World Series Win")]
    ax.legend(handles=legend_patches, fontsize=7.5, ncol=2, loc="upper left", frameon=True)
    fig.tight_layout(pad=1.5)
    return fig_bytes(fig)

# Radar
def normalize_col(series, invert=False):
    mn, mx = series.min(), series.max()
    if mx == mn: return pd.Series([50.0]*len(series), index=series.index)
    n = (series - mn)/(mx-mn)*100
    return 100-n if invert else n

def radar_chart(era_df):
    categories = ["Win %", "Run Diff", "OPS", "ERA inv", "FIP inv"]
    N = len(categories)
    angles = [n/float(N)*2*math.pi for n in range(N)]
    angles += angles[:1]

    era_df = era_df.copy()
    era_df["win_n"] = normalize_col(era_df["avg_wins"])
    era_df["rd_n"]  = normalize_col(era_df["avg_rd"])
    era_df["ops_n"] = normalize_col(era_df["avg_ops"])
    era_df["era_n"] = normalize_col(era_df["avg_era"], invert=True)
    era_df["fip_n"] = normalize_col(era_df["avg_fip"], invert=True)

    fig, ax = plt.subplots(figsize=(7, 7), subplot_kw=dict(polar=True), facecolor="white")
    ax.set_facecolor("#f8f8f8")
    for _, row in era_df.iterrows():
        vals = [row["win_n"], row["rd_n"], row["ops_n"], row["era_n"], row["fip_n"]]
        vals += vals[:1]
        c = hex_to_rgb_f(row["color"])
        ax.plot(angles, vals, color=c, linewidth=2.2)
        ax.fill(angles, vals, color=c, alpha=0.13)
        # label on peak
        peak_idx = vals.index(max(vals[:-1]))
        ax.annotate(row["short"], (angles[peak_idx], vals[peak_idx]),
                    textcoords="offset points", xytext=(8,4),
                    fontsize=7.5, color=c, fontweight="bold")

    ax.set_thetagrids([a*180/math.pi for a in angles[:-1]], categories, fontsize=9, color="#0d1b2a")
    ax.set_ylim(0,100)
    ax.set_yticklabels([])
    ax.grid(color="#ccc", linewidth=0.5)
    ax.set_title("Era Profile Comparison\n(normalised 0–100, higher = better)", fontsize=10,
                 fontweight="bold", color="#0d1b2a", pad=18)
    fig.tight_layout(pad=1.5)
    return fig_bytes(fig)

# Avg wins bar
def avg_wins_chart(era_df):
    fig, ax = plt.subplots(figsize=(10, 3.5), facecolor="white")
    names  = era_df["short"].tolist()
    vals   = era_df["avg_wins"].tolist()
    colors = [hex_to_rgb_f(c) for c in era_df["color"]]
    x = range(len(names))
    bars = ax.bar(x, vals, color=colors, width=0.6, zorder=3, edgecolor="white")
    ax.axhline(81, color="#c8102e", linestyle="--", linewidth=1.2, alpha=0.7, label=".500 (81 W)")
    for bar, v in zip(bars, vals):
        ax.text(bar.get_x()+bar.get_width()/2, v+0.5, f"{v}", ha="center", va="bottom",
                fontsize=9, fontweight="bold", color="#0d1b2a")
    ax.set_xticks(list(x)); ax.set_xticklabels(names, fontsize=8.5, rotation=10, ha="right")
    ax.set_ylim(60, 100)
    ax.legend(fontsize=8)
    styled_ax(ax, "Average Wins per Season by Era", ylabel="Avg Wins")
    fig.tight_layout(pad=1.2)
    return fig_bytes(fig)

# ══════════════════════════════════════════════════════════════════════════════
# PPTX LAYOUT HELPERS
# ══════════════════════════════════════════════════════════════════════════════
def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_picture_bytes(slide, img_bytes, left, top, width=None, height=None):
    if width and height:
        slide.shapes.add_picture(img_bytes, left, top, width=width, height=height)
    elif width:
        slide.shapes.add_picture(img_bytes, left, top, width=width)
    elif height:
        slide.shapes.add_picture(img_bytes, left, top, height=height)
    else:
        slide.shapes.add_picture(img_bytes, left, top)

def navy_header(slide, title, subtitle=None):
    """Full-width navy top bar with title."""
    add_rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    add_text(slide, title,
             Inches(0.35), Inches(0.10), Inches(12.5), Inches(0.65),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.35), Inches(0.72), Inches(12.5), Inches(0.38),
                 font_size=13, bold=False, color=GOLD, align=PP_ALIGN.LEFT)
    # Red accent line
    add_rect(slide, 0, Inches(1.15), SW, Pt(3), RED)

def kpi_box(slide, label, value, left, top, w, h, val_color=RED, bg=LGRAY):
    add_rect(slide, left, top, w, h, bg)
    add_text(slide, str(value), left, top+Pt(6), w, h*0.55,
             font_size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top+h*0.55, w, h*0.42,
             font_size=9, bold=False, color=DGRAY, align=PP_ALIGN.CENTER)

def era_chip(slide, text, color_hex, left, top):
    c = RGBColor(int(color_hex[1:3],16), int(color_hex[3:5],16), int(color_hex[5:7],16))
    add_rect(slide, left, top, Inches(1.75), Inches(0.32), c)
    add_text(slide, text, left+Pt(4), top+Pt(2), Inches(1.65), Inches(0.28),
             font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════
def slide_cover(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full navy background
    add_rect(slide, 0, 0, SW, SH, NAVY)
    # Red horizontal bands
    add_rect(slide, 0, Inches(2.5), SW, Pt(5), RED)
    add_rect(slide, 0, Inches(5.2), SW, Pt(5), RED)

    # Main title
    add_text(slide, "RED SOX FRANCHISE",
             Inches(0.5), Inches(0.6), Inches(12), Inches(1.0),
             font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "TIMELINE  1967 – 2024",
             Inches(0.5), Inches(1.55), Inches(12), Inches(0.85),
             font_size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "Seven Eras · 58 Seasons · 6 World Series Championships",
             Inches(0.5), Inches(2.62), Inches(12), Inches(0.55),
             font_size=16, bold=False, color=MGRAY, align=PP_ALIGN.CENTER)

    # KPI row
    ws   = int((df["playoff_result"]=="world series").sum())
    best = int(df["wins"].max())
    best_yr = int(df.loc[df["wins"].idxmax(),"season"])
    playoff = int(df["playoff_result"].isin(["world series","lost alcs","lost alds","al east tie-break loss"]).sum())
    kw, kh = Inches(2.9), Inches(1.55)
    ky = Inches(3.35)
    kpi_box(slide, "Total Seasons",       "58",           Inches(0.18), ky, kw, kh, WHITE, rgb(20,45,70))
    kpi_box(slide, "World Series Titles", str(ws),        Inches(3.22), ky, kw, kh, GOLD,  rgb(20,45,70))
    kpi_box(slide, "Best Season",         f"{best} W",    Inches(6.26), ky, kw, kh, WHITE, rgb(20,45,70))
    kpi_box(slide, "Best Year",           str(best_yr),   Inches(9.30), ky, kw, kh, GOLD,  rgb(20,45,70))

    add_text(slide,
             "From the Impossible Dream of '67 to the rebuild of today — "
             "the complete story of Boston Red Sox baseball.",
             Inches(1), Inches(5.2), Inches(11), Inches(0.65),
             font_size=13, bold=False, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)

    # Era chips row
    chip_x = Inches(0.35)
    for era, clr in ERA_COLORS_HEX.items():
        era_chip(slide, era.split("(")[0].strip(), clr, chip_x, Inches(6.1))
        chip_x += Inches(1.82)

def slide_story_in_numbers(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "The Story in Numbers", "58 seasons · every key franchise stat at a glance")

    stats = [
        ("58",   "Total Seasons"),
        ("6",    "World Series Titles"),
        ("108",  "Most Wins (2018)"),
        ("+229", "Best Run Diff (2018)"),
        ("20",   "Playoff Appearances"),
        ("10",   "Managers"),
        ("34%",  "Playoff Rate"),
        (f"{df['wins'].mean():.1f}", "Avg Wins / Season"),
    ]
    kw, kh = Inches(2.9), Inches(1.45)
    positions = [
        (Inches(0.18), Inches(1.35)), (Inches(3.22), Inches(1.35)),
        (Inches(6.26), Inches(1.35)), (Inches(9.30), Inches(1.35)),
        (Inches(0.18), Inches(3.0)),  (Inches(3.22), Inches(3.0)),
        (Inches(6.26), Inches(3.0)),  (Inches(9.30), Inches(3.0)),
    ]
    val_colors = [RED, GOLD, RED, GOLD, RED, GOLD, RED, GOLD]
    bg_colors  = [LGRAY]*8

    for (val, lbl), (lx, ly), vc in zip(stats, positions, val_colors):
        kpi_box(slide, lbl, val, lx, ly, kw, kh, vc)

    # World Series winners list
    add_rect(slide, Inches(0.18), Inches(4.62), SW-Inches(0.36), Inches(0.38), NAVY)
    add_text(slide, "  CHAMPIONSHIP YEARS:  " + "  ·  ".join(str(y) for y in sorted(
        df[df["playoff_result"]=="world series"]["season"].tolist())),
             Inches(0.22), Inches(4.64), SW-Inches(0.4), Inches(0.34),
             font_size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    # Bottom signal note
    add_text(slide,
             "Run Differential is used as the primary 'process' proxy. "
             "A team that wins more than its run diff predicts has likely overperformed; "
             "one that wins less has underperformed.",
             Inches(0.35), Inches(5.15), Inches(12.6), Inches(0.9),
             font_size=10, color=DGRAY, italic=True)

def slide_era_map(prs, df):
    """Era overview — coloured era band chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "The Seven Eras", "How franchise history breaks into distinct chapters")

    # Era table
    ew, eh = Inches(12.9), Inches(0.62)
    ey = Inches(1.25)
    for era, (yrs, clr) in MACRO_ERAS.items():
        sub = df[df["season"].isin(yrs)]
        if len(sub) == 0: continue
        avg_w  = sub["wins"].mean()
        ws_n   = (sub["playoff_result"]=="world series").sum()
        tagline = ERA_TAGLINES[era]

        # Row bg
        c = RGBColor(int(clr[1:3],16), int(clr[3:5],16), int(clr[5:7],16))
        add_rect(slide, Inches(0.22), ey, Inches(1.7), eh, c)
        add_rect(slide, Inches(1.95), ey, Inches(11.17), eh, LGRAY)

        add_text(slide, era, Inches(0.25), ey+Pt(4), Inches(1.65), eh-Pt(8),
                 font_size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide,
                 f"{min(yrs)}–{max(yrs)}   |   Avg {avg_w:.0f} W   |   {'★ '*ws_n if ws_n else '—'}   |   {tagline}",
                 Inches(2.0), ey+Pt(4), Inches(11.0), eh-Pt(8),
                 font_size=9.5, bold=False, color=NAVY, align=PP_ALIGN.LEFT)
        ey += eh + Pt(4)

    add_text(slide, "★ = World Series Championship",
             Inches(0.22), Inches(6.85), Inches(5), Inches(0.35),
             font_size=8.5, color=DGRAY, italic=True)

def slide_era_deepdive(prs, df, era_name, yrs, color_hex):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)

    # Coloured left rail
    c = RGBColor(int(color_hex[1:3],16), int(color_hex[3:5],16), int(color_hex[5:7],16))
    add_rect(slide, 0, 0, Inches(0.22), SH, c)

    # Header
    add_rect(slide, Inches(0.22), 0, SW-Inches(0.22), Inches(1.15), NAVY)
    add_text(slide, era_name,
             Inches(0.4), Inches(0.08), Inches(9), Inches(0.65),
             font_size=24, bold=True, color=WHITE)
    add_text(slide, ERA_TAGLINES[era_name],
             Inches(0.4), Inches(0.72), Inches(9), Inches(0.4),
             font_size=12, color=GOLD)
    add_rect(slide, Inches(0.22), Inches(1.15), SW-Inches(0.22), Pt(3), RED)

    sub = df[df["season"].isin(yrs)].copy()
    sub["tag"] = sub.apply(classify_process_vs_result, axis=1)

    # KPI boxes top-right
    ws_n   = int((sub["playoff_result"]=="world series").sum())
    avg_w  = sub["wins"].mean()
    best_s = int(sub.loc[sub["wins"].idxmax(),"season"])
    best_w = int(sub["wins"].max())
    avg_rd = sub["run_diff"].mean()

    add_rect(slide, Inches(9.5), Inches(1.25), Inches(3.6), Inches(1.25), LGRAY)
    add_text(slide, f"{avg_w:.0f}", Inches(9.55), Inches(1.28), Inches(1.1), Inches(0.65),
             font_size=26, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(slide, "Avg Wins", Inches(9.55), Inches(1.88), Inches(1.1), Inches(0.35),
             font_size=8.5, color=DGRAY, align=PP_ALIGN.CENTER)
    add_text(slide, f"{avg_rd:+.0f}", Inches(10.75), Inches(1.28), Inches(1.1), Inches(0.65),
             font_size=26, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "Avg Run Diff", Inches(10.75), Inches(1.88), Inches(1.1), Inches(0.35),
             font_size=8.5, color=DGRAY, align=PP_ALIGN.CENTER)
    ws_val = "★"*ws_n if ws_n else "—"
    add_text(slide, ws_val, Inches(11.95), Inches(1.28), Inches(1.0), Inches(0.65),
             font_size=22 if ws_n else 18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, "WS Titles", Inches(11.95), Inches(1.88), Inches(1.0), Inches(0.35),
             font_size=8.5, color=DGRAY, align=PP_ALIGN.CENTER)

    # Era wins bar chart (left side, most of width)
    chart = era_wins_chart(sub, era_name, color_hex)
    add_picture_bytes(slide, chart, Inches(0.3), Inches(1.28), width=Inches(8.9))

    # Key players
    add_rect(slide, Inches(0.3), Inches(4.55), Inches(12.7), Inches(0.32), c)
    add_text(slide, "  KEY PLAYERS", Inches(0.35), Inches(4.57), Inches(3.5), Inches(0.28),
             font_size=9, bold=True, color=WHITE)
    add_text(slide, ERA_KEY_PLAYERS[era_name],
             Inches(3.8), Inches(4.57), Inches(9.2), Inches(0.28),
             font_size=9, color=WHITE)

    # Seasons table
    col_header = ["Year","W","L","RD","Result","OPS","ERA","Manager"]
    col_widths  = [0.7, 0.45, 0.45, 0.55, 1.45, 0.65, 0.6, 1.65]
    row_h = Inches(0.275)
    ty = Inches(4.97)
    tx_starts = [Inches(0.3)]
    for w in col_widths[:-1]:
        tx_starts.append(tx_starts[-1]+Inches(w))

    # Header row
    add_rect(slide, Inches(0.3), ty, Inches(sum(col_widths)), row_h, NAVY)
    for header, lx, cw in zip(col_header, tx_starts, col_widths):
        add_text(slide, header, lx+Pt(2), ty+Pt(3), Inches(cw)-Pt(4), row_h-Pt(6),
                 font_size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ty += row_h

    for idx, (_, row) in enumerate(sub.sort_values("season").iterrows()):
        bg = LGRAY if idx%2==0 else WHITE
        if row["playoff_result"] == "world series":
            bg = rgb(255, 248, 220)   # gold tint
        add_rect(slide, Inches(0.3), ty, Inches(sum(col_widths)), row_h, bg)
        result_str = str(row["playoff_result"]).replace("world series","WS WIN").replace(
            "lost alcs","ALCS").replace("lost alds","ALDS").replace(
            "al east tie-break loss","TBL").replace("missed playoffs","—")
        cells = [str(int(row["season"])), str(int(row["wins"])), str(int(row["losses"])),
                 f"{row['run_diff']:+d}", result_str,
                 f"{row['team_ops']:.3f}", f"{row['team_era']:.2f}", str(row["manager"])]
        for val, lx, cw in zip(cells, tx_starts, col_widths):
            text_color = NAVY
            if row["playoff_result"] == "world series" and val == "WS WIN":
                text_color = RED
            add_text(slide, val, lx+Pt(2), ty+Pt(2), Inches(cw)-Pt(4), row_h-Pt(4),
                     font_size=7.5, color=text_color, align=PP_ALIGN.CENTER)
        ty += row_h
        if ty > Inches(7.3): break

def slide_ws_championships(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "The Six Championships", "Every World Series title in franchise history")

    ws_rows = df[df["playoff_result"]=="world series"].sort_values("season")
    card_w = Inches(2.05)
    card_h = Inches(5.55)
    card_x = Inches(0.18)
    card_y = Inches(1.25)

    for _, row in ws_rows.iterrows():
        era_color = "#4a90d9"
        for era, (yrs, clr) in MACRO_ERAS.items():
            if int(row["season"]) in yrs:
                era_color = clr; break
        c = RGBColor(int(era_color[1:3],16), int(era_color[3:5],16), int(era_color[5:7],16))

        # Card bg
        add_rect(slide, card_x, card_y, card_w, card_h, LGRAY)
        # Top color strip
        add_rect(slide, card_x, card_y, card_w, Inches(0.55), c)
        # Year
        add_text(slide, str(int(row["season"])),
                 card_x, card_y+Pt(3), card_w, Inches(0.5),
                 font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Trophy
        add_text(slide, "WORLD SERIES",
                 card_x, card_y+Inches(0.58), card_w, Inches(0.32),
                 font_size=8, bold=True, color=RED, align=PP_ALIGN.CENTER)
        # Record
        add_text(slide, f"{int(row['wins'])}-{int(row['losses'])}",
                 card_x, card_y+Inches(0.92), card_w, Inches(0.42),
                 font_size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, "W-L",
                 card_x, card_y+Inches(1.28), card_w, Inches(0.28),
                 font_size=8, color=DGRAY, align=PP_ALIGN.CENTER)
        # Run diff
        add_text(slide, f"+{int(row['run_diff'])}",
                 card_x, card_y+Inches(1.6), card_w, Inches(0.38),
                 font_size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        add_text(slide, "Run Diff",
                 card_x, card_y+Inches(1.9), card_w, Inches(0.28),
                 font_size=8, color=DGRAY, align=PP_ALIGN.CENTER)
        # Manager
        add_text(slide, row["manager"],
                 card_x, card_y+Inches(2.25), card_w, Inches(0.35),
                 font_size=8.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(slide, "Manager",
                 card_x, card_y+Inches(2.55), card_w, Inches(0.28),
                 font_size=7.5, color=DGRAY, align=PP_ALIGN.CENTER)
        # Key players
        players = str(row["key_players"])
        wrapped = textwrap.fill(players, width=22)
        add_text(slide, wrapped,
                 card_x+Pt(4), card_y+Inches(2.88), card_w-Pt(8), Inches(1.4),
                 font_size=7.5, color=DGRAY, wrap=True, align=PP_ALIGN.CENTER)
        # Mechanism
        mech = str(row["mechanism_summary"])
        if len(mech) > 90: mech = mech[:88]+"…"
        add_text(slide, mech,
                 card_x+Pt(4), card_y+Inches(4.32), card_w-Pt(8), Inches(1.1),
                 font_size=6.5, color=DGRAY, wrap=True, align=PP_ALIGN.CENTER, italic=True)

        card_x += card_w + Inches(0.08)

def slide_process_vs_result(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "Process vs Result", "Run Differential as a proxy for true team quality")

    # Scatter chart
    sc = scatter_chart(df)
    add_picture_bytes(slide, sc, Inches(0.25), Inches(1.25), width=Inches(8.4))

    # Explainer panel right side
    add_rect(slide, Inches(8.75), Inches(1.25), Inches(4.35), Inches(6.15), LGRAY)
    add_text(slide, "How to read this chart",
             Inches(8.85), Inches(1.35), Inches(4.15), Inches(0.42),
             font_size=11, bold=True, color=NAVY)
    bullets = [
        ("X-axis:", "Run Differential — a better\npredictor of true quality than wins"),
        ("Y-axis:", "Actual Win Total"),
        ("Top-right:", "High wins + high RD\n→ Genuine dominance"),
        ("Top-left:", "High wins, low RD\n→ Lucky / overperformed"),
        ("Gold stars:", "World Series champions"),
        ("",""),
        ("Signal seasons:", "Wins backed by elite\nrun differential"),
        ("Heartbreak:", "1978 (99W) — best RD\nthat didn't win the WS"),
        ("Best process:", "2018 (108W, +229 RD)"),
    ]
    by = Inches(1.85)
    for label, body in bullets:
        if not label and not body:
            by += Inches(0.12); continue
        add_text(slide, label, Inches(8.9), by, Inches(1.5), Inches(0.5),
                 font_size=8.5, bold=True, color=RED)
        add_text(slide, body, Inches(10.4), by, Inches(2.55), Inches(0.5),
                 font_size=8.5, color=NAVY)
        by += Inches(0.52)

def slide_wins_timeline(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "58 Seasons at a Glance", "Every win total 1967–2024 coloured by era")

    chart = franchise_timeline_chart(df)
    add_picture_bytes(slide, chart, Inches(0.18), Inches(1.28), width=Inches(12.95))

    # Era band legend below
    lx = Inches(0.35)
    for era, clr in ERA_COLORS_HEX.items():
        era_chip(slide, era.split("(")[0].strip(), clr, lx, Inches(6.42))
        lx += Inches(1.82)

def slide_era_radar(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "Era Profile Radar", "Normalised 0–100 · higher = better on every axis")

    # Build era_df
    era_records = []
    era_map = {y: era for era, (yrs,_) in MACRO_ERAS.items() for y in yrs}
    df2 = df.copy(); df2["macro_era"] = df2["season"].map(era_map)
    for era, (yrs, clr) in MACRO_ERAS.items():
        sub = df2[df2["macro_era"]==era]
        era_records.append({
            "era": era, "short": era.split("(")[0].strip(), "color": clr,
            "avg_wins": sub["wins"].mean(), "avg_rd": sub["run_diff"].mean(),
            "avg_ops": sub["team_ops"].mean(), "avg_era": sub["team_era"].mean(),
            "avg_fip": sub["team_fip"].mean(),
        })
    era_df = pd.DataFrame(era_records)

    radar = radar_chart(era_df)
    add_picture_bytes(slide, radar, Inches(0.2), Inches(1.25), height=Inches(6.0))

    # Insights panel
    add_rect(slide, Inches(7.0), Inches(1.25), Inches(6.15), Inches(6.05), LGRAY)
    insights = [
        ("Francona Dynasty", "Leads Run Diff and Win % — the most complete era"),
        ("Pedro & Manny",    "Best OPS axis — peak offensive production"),
        ("Yaz Era",          "Best ERA axis — dominated a pitcher-friendly era"),
        ("Betts Emergence",  "Strong balanced profile — 2018 was the peak"),
        ("Rebuild & Now",    "Weakest overall — negative run diff since 2020"),
        ("Boggs & Clemens",  "Win % limits despite strong pitching — no title"),
        ("Rice & Lynn",      "Power-era offense but inconsistent pitching"),
    ]
    add_text(slide, "Key Takeaways",
             Inches(7.15), Inches(1.35), Inches(5.85), Inches(0.42),
             font_size=12, bold=True, color=NAVY)
    iy = Inches(1.85)
    for era_short, note in insights:
        add_text(slide, f"• {era_short}", Inches(7.2), iy, Inches(2.4), Inches(0.5),
                 font_size=9, bold=True, color=RED)
        add_text(slide, note, Inches(9.6), iy, Inches(3.4), Inches(0.5),
                 font_size=9, color=NAVY)
        iy += Inches(0.55)

def slide_key_players(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "Key Players by Era", "The names that defined each chapter")

    ey = Inches(1.28)
    for era, (yrs, clr) in MACRO_ERAS.items():
        c = RGBColor(int(clr[1:3],16), int(clr[3:5],16), int(clr[5:7],16))
        sub = df[df["season"].isin(yrs)]
        ws_count = int((sub["playoff_result"]=="world series").sum())
        ws_label = ("  ★"*ws_count) if ws_count else ""

        add_rect(slide, Inches(0.22), ey, Inches(2.5), Inches(0.62), c)
        add_text(slide, era+ws_label,
                 Inches(0.27), ey+Pt(4), Inches(2.4), Inches(0.54),
                 font_size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        add_rect(slide, Inches(2.75), ey, Inches(10.4), Inches(0.62), LGRAY)
        add_text(slide, ERA_KEY_PLAYERS[era],
                 Inches(2.82), ey+Pt(6), Inches(10.2), Inches(0.5),
                 font_size=9.5, color=NAVY, align=PP_ALIGN.LEFT)
        ey += Inches(0.66)

    add_text(slide, "★ = World Series championship earned in that era",
             Inches(0.22), Inches(6.85), Inches(6), Inches(0.38),
             font_size=8, color=DGRAY, italic=True)

def slide_current_state(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, WHITE)
    navy_header(slide, "Current State & What's Next", "2020–2024 · The Rebuild & Now Era")

    sub = df[df["season"].isin(range(2020,2025))].copy()
    sub["tag"] = sub.apply(classify_process_vs_result, axis=1)

    # Left: seasons table
    col_h = ["Year","W","L","RD","Result","OPS","ERA","Manager","Key Players"]
    col_w  = [0.58, 0.42, 0.42, 0.52, 1.1, 0.6, 0.55, 1.35, 2.15]
    tx_starts = [Inches(0.28)]
    for w in col_w[:-1]: tx_starts.append(tx_starts[-1]+Inches(w))
    row_h = Inches(0.42)
    ty = Inches(1.3)
    total_w = Inches(sum(col_w))

    add_rect(slide, Inches(0.28), ty, total_w, row_h, NAVY)
    for header, lx, cw in zip(col_h, tx_starts, col_w):
        add_text(slide, header, lx+Pt(2), ty+Pt(4), Inches(cw)-Pt(4), row_h-Pt(8),
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    ty += row_h

    for idx, (_, row) in enumerate(sub.sort_values("season").iterrows()):
        bg = LGRAY if idx%2==0 else WHITE
        add_rect(slide, Inches(0.28), ty, total_w, row_h, bg)
        result_str = str(row["playoff_result"]).replace("world series","WS").replace(
            "lost alcs","ALCS").replace("lost alds","ALDS").replace(
            "al east tie-break loss","TBL").replace("missed playoffs","—")
        players = str(row["key_players"])
        cells = [str(int(row["season"])), str(int(row["wins"])), str(int(row["losses"])),
                 f"{row['run_diff']:+d}", result_str,
                 f"{row['team_ops']:.3f}", f"{row['team_era']:.2f}",
                 str(row["manager"]), players]
        for val, lx, cw in zip(cells, tx_starts, col_w):
            add_text(slide, val, lx+Pt(2), ty+Pt(3), Inches(cw)-Pt(4), row_h-Pt(6),
                     font_size=8, color=NAVY, align=PP_ALIGN.CENTER)
        ty += row_h

    # Right: outlook panel
    add_rect(slide, Inches(7.55), Inches(1.28), Inches(5.55), Inches(5.9), rgb(20,45,70))
    add_text(slide, "The Road Ahead",
             Inches(7.7), Inches(1.38), Inches(5.2), Inches(0.5),
             font_size=14, bold=True, color=WHITE)
    add_rect(slide, Inches(7.55), Inches(1.9), Inches(5.55), Pt(2), RED)

    points = [
        ("Rafael Devers", "Franchise cornerstone — elite bat, needs contract security"),
        ("Jarren Duran",  "2023–24 breakout — the process engine of the current roster"),
        ("Brayan Bello",  "Ace candidate — FIP better than ERA in 2024"),
        ("Tanner Houck",  "Frontline starter emerging — key to 2025 rotation"),
        ("Farm system",   "Ranked mid-tier — international pipeline developing"),
        ("2025 outlook",  "81 W baseline · needs rotation depth and bullpen upgrades"),
    ]
    py = Inches(2.05)
    for player, note in points:
        add_text(slide, player, Inches(7.7), py, Inches(2.1), Inches(0.46),
                 font_size=9.5, bold=True, color=GOLD)
        add_text(slide, note, Inches(9.82), py, Inches(3.1), Inches(0.46),
                 font_size=9, color=MGRAY)
        py += Inches(0.52)

    add_rect(slide, Inches(7.55), Inches(5.2), Inches(5.55), Inches(0.98), RED)
    add_text(slide, '"The foundation is there. The question is timing."',
             Inches(7.7), Inches(5.28), Inches(5.2), Inches(0.78),
             font_size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SW, SH, NAVY)
    add_rect(slide, 0, Inches(3.1), SW, Pt(5), RED)
    add_rect(slide, 0, Inches(4.7), SW, Pt(5), RED)

    add_text(slide, "RED SOX FRANCHISE TIMELINE",
             Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.9),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "1967 – 2024",
             Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.7),
             font_size=28, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, "6 Championships · 58 Seasons · One City's Obsession",
             Inches(0.5), Inches(2.42), Inches(12.3), Inches(0.52),
             font_size=15, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)

    add_text(slide, "1967  ·  1975  ·  1978  ·  1986  ·  2004  ·  2007  ·  2013  ·  2018",
             Inches(0.5), Inches(3.28), Inches(12.3), Inches(0.48),
             font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, "Built with frame2-redsox-history analytics engine · redsox_history.csv · 2025",
             Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.45),
             font_size=9, color=DGRAY, align=PP_ALIGN.CENTER)

    # Era chips
    chip_x = Inches(0.35)
    for era, clr in ERA_COLORS_HEX.items():
        era_chip(slide, era.split("(")[0].strip(), clr, chip_x, Inches(5.05))
        chip_x += Inches(1.82)

# ══════════════════════════════════════════════════════════════════════════════
# MAIN BUILD
# ══════════════════════════════════════════════════════════════════════════════
def build_slideshow(output_path):
    print("Loading data…")
    df = load_redsox_history()
    df["season"] = df["season"].astype(int)
    df["tag"] = df.apply(classify_process_vs_result, axis=1)

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    print("Slide 01/17 — Cover…")
    slide_cover(prs, df)

    print("Slide 02/17 — Story in Numbers…")
    slide_story_in_numbers(prs, df)

    print("Slide 03/17 — Era Map…")
    slide_era_map(prs, df)

    era_list = list(MACRO_ERAS.items())
    for i, (era_name, (yrs, clr)) in enumerate(era_list):
        print(f"Slide {i+4:02d}/17 — Era deep-dive: {era_name}…")
        slide_era_deepdive(prs, df, era_name, yrs, clr)

    print("Slide 11/17 — Championships…")
    slide_ws_championships(prs, df)

    print("Slide 12/17 — Process vs Result…")
    slide_process_vs_result(prs, df)

    print("Slide 13/17 — Wins Timeline…")
    slide_wins_timeline(prs, df)

    print("Slide 14/17 — Era Radar…")
    slide_era_radar(prs, df)

    print("Slide 15/17 — Key Players…")
    slide_key_players(prs, df)

    print("Slide 16/17 — Current State…")
    slide_current_state(prs, df)

    print("Slide 17/17 — Closing…")
    slide_closing(prs)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"\n✓ Saved {len(prs.slides)} slides → {output_path}  ({size_kb} KB)")
    return output_path

if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "redsox_franchise_timeline.pptx"
    build_slideshow(out)
